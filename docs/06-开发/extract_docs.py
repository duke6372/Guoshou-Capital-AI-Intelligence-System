from __future__ import annotations

import argparse
import json
import os
import re
import shutil
import subprocess
from dataclasses import asdict, dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable


@dataclass(frozen=True)
class ExtractRecord:
    source_path: str
    file_type: str
    extracted_at: str
    extractor: str
    pages_or_slides: int
    total_chars: int
    avg_chars: int
    empty_units: int
    quality: str
    output_md: str


def _utc_now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def _readable_slug(name: str) -> str:
    name = Path(name).stem
    name = re.sub(r"\s+", " ", name).strip()
    name = re.sub(r"[\/\\:*?\"<>|]+", "_", name)
    return name


def _quality(total_chars: int, unit_count: int, empty_units: int) -> str:
    if unit_count <= 0:
        return "low"
    avg = total_chars // unit_count
    empty_ratio = empty_units / max(unit_count, 1)
    if avg >= 800 and empty_ratio <= 0.15:
        return "high"
    if avg >= 200 and empty_ratio <= 0.4:
        return "medium"
    return "low"


def _write_md(output_path: Path, header: dict[str, Any], sections: list[tuple[str, str]]) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    lines: list[str] = []
    lines.append("---")
    for k, v in header.items():
        lines.append(f"{k}: {json.dumps(v, ensure_ascii=False)}")
    lines.append("---")
    lines.append("")
    for title, body in sections:
        lines.append(f"## {title}")
        lines.append("")
        lines.append(body.strip() if body.strip() else "(空)")
        lines.append("")
    output_path.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")


def _extract_pdf_with_pdfplumber(source_path: Path) -> tuple[str, list[tuple[str, str]]]:
    import pdfplumber  # type: ignore

    sections: list[tuple[str, str]] = []
    with pdfplumber.open(str(source_path)) as pdf:
        for idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            sections.append((f"Page {idx}", text))
    return "pdfplumber", sections


def _extract_pdf_with_pypdf(source_path: Path) -> tuple[str, list[tuple[str, str]]]:
    from pypdf import PdfReader  # type: ignore

    reader = PdfReader(str(source_path))
    sections: list[tuple[str, str]] = []
    for idx, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        sections.append((f"Page {idx}", text))
    return "pypdf", sections


def _extract_pdf_with_pdftotext(source_path: Path) -> tuple[str, list[tuple[str, str]]]:
    if shutil.which("pdftotext") is None:
        raise RuntimeError("pdftotext not found")
    proc = subprocess.run(
        ["pdftotext", "-layout", "-nopgbrk", str(source_path), "-"],
        check=False,
        capture_output=True,
        text=True,
    )
    text = (proc.stdout or "").strip()
    return "pdftotext", [("Text", text)]


def _extract_pdf(source_path: Path) -> tuple[str, list[tuple[str, str]]]:
    try:
        return _extract_pdf_with_pdfplumber(source_path)
    except Exception:
        pass
    try:
        return _extract_pdf_with_pypdf(source_path)
    except Exception:
        pass
    return _extract_pdf_with_pdftotext(source_path)


def _extract_pptx(source_path: Path) -> tuple[str, list[tuple[str, str]]]:
    from pptx import Presentation  # type: ignore

    pres = Presentation(str(source_path))
    sections: list[tuple[str, str]] = []
    for idx, slide in enumerate(pres.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if isinstance(text, str) and text.strip():
                parts.append(text.strip())
        notes_text = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes_text = ""
        if notes_text:
            parts.append("Notes:\n" + notes_text)
        sections.append((f"Slide {idx}", "\n\n".join(parts)))
    return "python-pptx", sections


def _iter_inputs(input_dir: Path) -> Iterable[Path]:
    for p in sorted(input_dir.glob("*")):
        if p.is_file() and p.suffix.lower() in {".pdf", ".pptx"}:
            yield p


def extract_one(source_path: Path, output_dir: Path) -> ExtractRecord:
    file_type = source_path.suffix.lower().lstrip(".")
    slug = _readable_slug(source_path.name)
    output_md = output_dir / f"{slug}.{file_type}.md"

    if file_type == "pdf":
        extractor, sections = _extract_pdf(source_path)
    elif file_type == "pptx":
        extractor, sections = _extract_pptx(source_path)
    else:
        raise ValueError(f"unsupported type: {file_type}")

    unit_count = len(sections)
    total_chars = sum(len(body or "") for _, body in sections)
    empty_units = sum(1 for _, body in sections if not (body or "").strip())
    quality = _quality(total_chars=total_chars, unit_count=unit_count, empty_units=empty_units)

    header = {
        "source_path": str(source_path),
        "file_type": file_type,
        "extracted_at": _utc_now_iso(),
        "extractor": extractor,
        "units": unit_count,
        "total_chars": total_chars,
        "avg_chars": (total_chars // unit_count) if unit_count else 0,
        "empty_units": empty_units,
        "quality": quality,
    }
    _write_md(output_md, header=header, sections=sections)

    return ExtractRecord(
        source_path=str(source_path),
        file_type=file_type,
        extracted_at=header["extracted_at"],
        extractor=extractor,
        pages_or_slides=unit_count,
        total_chars=total_chars,
        avg_chars=header["avg_chars"],
        empty_units=empty_units,
        quality=quality,
        output_md=str(output_md),
    )


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", default="Project Files")
    parser.add_argument("--output", default="docs/source")
    args = parser.parse_args()

    repo_root = Path(os.getcwd())
    input_dir = (repo_root / args.input).resolve()
    output_dir = (repo_root / args.output).resolve()

    if not input_dir.exists():
        raise SystemExit(f"input directory not found: {input_dir}")

    output_dir.mkdir(parents=True, exist_ok=True)

    records: list[ExtractRecord] = []
    for source_path in _iter_inputs(input_dir):
        record = extract_one(source_path=source_path, output_dir=output_dir)
        records.append(record)

    manifest = {
        "generated_at": _utc_now_iso(),
        "input_dir": str(input_dir),
        "output_dir": str(output_dir),
        "records": [asdict(r) for r in records],
    }
    (output_dir / "manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

